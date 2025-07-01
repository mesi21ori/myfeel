from flask import Blueprint, request, jsonify
import os
import docx2txt
import fitz  
from pptx import Presentation
from PIL import Image
import pytesseract
import google.generativeai as genai
from config import Config

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

analyze_bp = Blueprint("analyze", __name__)
genai.configure(api_key=Config.GOOGLE_API_KEY)

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.docx':
        return docx2txt.process(file_path)
    
    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    elif ext == '.pdf':
        text = ""
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        return text
    
    elif ext == '.pptx':
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif ext in ['.jpg', '.jpeg', '.png']:
        try:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        except Exception as e:
            return f"Error extracting text from image: {str(e)}"

    else:
        return ""

@analyze_bp.route('/', methods=['POST'])
def analyze_document():
    data = request.get_json()
    file_path = data.get('path')
    analysis_type = data.get('type')
    user_message = data.get('message', '')

    if not file_path or not analysis_type:
        return jsonify({'error': 'Missing "path" or "type"'}), 400

    file_text = extract_text(file_path).strip()
    if not file_text:
        return jsonify({'error': 'No text found in the document'}), 400

    prompt = f"{user_message.strip()}\n\nDocument:\n{file_text}" if user_message else f"Analyze this document and provide a {analysis_type}:\n\n{file_text}"

    try:
        model = genai.GenerativeModel("models/gemini-1.5-flash")  
        response = model.generate_content(prompt)
        return jsonify({
            'type': analysis_type,
            'result': response.text.strip()
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


